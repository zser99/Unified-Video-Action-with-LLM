#!/usr/bin/env python3
"""UVA to CoT 기획서 PPTX 생성"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

C_DARK   = RGBColor(0x7B, 0x32, 0x00)
C_FOOTER = RGBColor(0x8B, 0x3A, 0x00)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BODY   = RGBColor(0x1A, 0x1A, 0x1A)
C_GRAY   = RGBColor(0x88, 0x88, 0x88)
BLANK    = prs.slide_layouts[6]
FONT_KO  = "맑은 고딕"
FONT_EN  = "Calibri"


# ── helpers ────────────────────────────────────────────────
def apply_gradient(shape, c1="E07000", c2="D4B200"):
    sp   = shape._element
    spPr = sp.find(qn("p:spPr"))
    keep = {qn("a:xfrm"), qn("a:prstGeom"), qn("a:custGeom"), qn("a:ln")}
    for child in list(spPr):
        if child.tag not in keep:
            spPr.remove(child)
    ns  = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    xml = (f'<a:gradFill {ns}><a:gsLst>'
           f'<a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>'
           f'<a:gs pos="100000"><a:srgbClr val="{c2}"/></a:gs>'
           f'</a:gsLst><a:lin ang="0" scaled="0"/></a:gradFill>')
    idx = 0
    for i, ch in enumerate(spPr):
        if ch.tag in (qn("a:xfrm"), qn("a:prstGeom")):
            idx = i + 1
    spPr.insert(idx, etree.fromstring(xml))


def add_hdr(slide, title, num=None):
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.8))
    bar.line.fill.background()
    apply_gradient(bar)
    tf = bar.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.2)
    tf.margin_top  = Inches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(22)
    r.font.name = FONT_KO
    r.font.color.rgb = C_WHITE
    if num:
        nb = slide.shapes.add_textbox(W - Inches(0.55), Inches(0.07),
                                      Inches(0.45), Inches(0.65))
        nb.text_frame.margin_top = 0
        p2 = nb.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = str(num)
        r2.font.bold  = True
        r2.font.size  = Pt(20)
        r2.font.name  = FONT_EN
        r2.font.color.rgb = C_WHITE


def add_ftr(slide):
    fb = slide.shapes.add_textbox(Inches(0.15), H - Inches(0.4),
                                  Inches(2.5), Inches(0.35))
    p = fb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Dongguk Univ."
    r.font.bold = True
    r.font.size = Pt(9)
    r.font.name = FONT_EN
    r.font.color.rgb = C_FOOTER


def new_tf(slide, top=Inches(0.85), left=Inches(0.35), w=None, h=None):
    if w is None: w = W - Inches(0.7)
    if h is None: h = H - top - Inches(0.42)
    box = slide.shapes.add_textbox(left, top, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    return tf


def _set_marL(p, inches):
    """Set paragraph left margin via XML (paragraph_format not in pptx 1.x)"""
    pPr = p._p.find(qn("a:pPr"))
    if pPr is None:
        pPr = etree.SubElement(p._p, qn("a:pPr"))
        p._p.insert(0, pPr)
    pPr.set("marL", str(int(Inches(inches))))


def h1(tf, text, sp=5):
    p = tf.add_paragraph()
    p.space_before = Pt(sp)
    r = p.add_run()
    r.text = "•  " + text
    r.font.bold  = True
    r.font.size  = Pt(13.5)
    r.font.name  = FONT_KO
    r.font.color.rgb = C_DARK
    return p


def li(tf, text, size=Pt(12), sp=2):
    p = tf.add_paragraph()
    p.space_before = Pt(sp)
    _set_marL(p, 0.38)
    r = p.add_run()
    r.text = "–  " + text
    r.font.size  = size
    r.font.name  = FONT_KO
    r.font.color.rgb = C_BODY
    return p


def ph_box(slide, left, top, w, h, label):
    """Gray placeholder box with centered label"""
    box = slide.shapes.add_textbox(left, top, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(10)
    r.font.name = FONT_KO
    r.font.color.rgb = C_GRAY
    return box


# ══════════════════════════════════════════════════════════
# Slide 1 ── 표지
# ══════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)

top_bar = s1.shapes.add_shape(1, 0, 0, W, Inches(1.05))
top_bar.line.fill.background()
apply_gradient(top_bar)

tb = s1.shapes.add_textbox(Inches(0.25), Inches(0.18),
                            W - Inches(0.5), Inches(0.7))
p = tb.text_frame.paragraphs[0]
r = p.add_run()
r.text = "Topic:  편하신대로 작성하시면되요"
r.font.bold = True; r.font.size = Pt(15)
r.font.name = FONT_KO; r.font.color.rgb = C_DARK

ct = s1.shapes.add_textbox(Inches(1.2), Inches(2.6),
                            W - Inches(2.4), Inches(1.8))
ct.text_frame.word_wrap = True
p2 = ct.text_frame.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "UVA to CoT 기획서"
r2.font.bold = True; r2.font.size = Pt(36)
r2.font.name = FONT_KO; r2.font.color.rgb = C_DARK

tb2 = s1.shapes.add_textbox(Inches(0.5), H - Inches(2.1),
                              Inches(5), Inches(0.5))
p3 = tb2.text_frame.paragraphs[0]
r3 = p3.add_run()
r3.text = "Presented by:"
r3.font.bold = True; r3.font.size = Pt(14)
r3.font.name = FONT_KO; r3.font.color.rgb = C_DARK
add_ftr(s1)


# ══════════════════════════════════════════════════════════
# Slide 2 ── 0. 기존 모델
# ══════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_hdr(s2, "0. 기존 모델", 2); add_ftr(s2)

tf2a = new_tf(s2, top=Inches(0.82), h=Inches(1.0))
h1(tf2a, "Backbone", sp=2)
li(tf2a, "UVA"); li(tf2a, "CoT")

ph_box(s2, Inches(0.5), Inches(1.95), W - Inches(1.0), Inches(3.0),
       "[ UVA architecture 이미지 ]                    [ CoT architecture 이미지 ]")

lbl = s2.shapes.add_textbox(Inches(0.5), Inches(5.03), W - Inches(1.0), Inches(0.35))
p_l = lbl.text_frame.paragraphs[0]; p_l.alignment = PP_ALIGN.CENTER
r_l = p_l.add_run()
r_l.text = "UVA architecture                                   CoT architecture"
r_l.font.size = Pt(10); r_l.font.name = FONT_KO; r_l.font.color.rgb = C_BODY

tf2b = new_tf(s2, top=Inches(5.45), h=Inches(1.5))
h1(tf2b, "간단 소개", sp=2)
li(tf2b, "발표자(팀)와 주제를 1-2 문장으로 소개. "
         "위 BaseModel 섹션 아래 이미지 영역에 두 모델의 architecture 그림을 첨부",
   size=Pt(11))


# ══════════════════════════════════════════════════════════
# Slide 3 ── 1. 연구 배경 및 동기
# ══════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_hdr(s3, "1. 연구 배경 및 동기", 3); add_ftr(s3)
tf3 = new_tf(s3)

h1(tf3, "UVA 모델의 한계점", sp=2)
li(tf3, "블랙박스 문제 : long-Horizon 태스크에서 명시적 하위목표, 대안 액션 비교에 대한 "
        "설명이 모자람. 또한 추론 시점에 추론 단계를 파악하기가 어렵다.")

h1(tf3, "UVA의 기여")
li(tf3, "비디오와 액션을 통합 잠재공간에서 학습하여 시각-행동 동역학을 효과적으로 결합.")
li(tf3, "디코더 분리로 정책 추론 시 비디오 생성을 생략해 속도 확보.")
li(tf3, "마스킹 기반 학습으로 policy/forward/inverse dynamics 등 다방면 학습이 가능함. "
        "이는 robustness 향상과 과적합을 방지함.")

h1(tf3, "CoT(Chain-of-Thought) 추론의 가능성")
li(tf3, "CoT는 정형화 된 추론을 강제하며, 중간 추론 결과를 명시적인 토큰으로 생성함으로써 "
        "복잡한 문제의 장기 의존성을 다루는 데 유리하다.")
li(tf3, "로봇 제어에서는 고수준 목표 → 중간 목표 → 저수준 행동의 계층화를 통해 "
        "오류 누적을 줄이고 재계획을 쉽게 만든다.")
li(tf3, "추론 과정을 명시화하면 디버깅/해석성도 개선된다.")

h1(tf3, "두 방법을 결합해야 하는 이유")
li(tf3, "CoT는 계획·비교·설명을, UVA는 학습된 동역학/정책을 담당한다.")
li(tf3, "UVA를 동결(frozen)하여 훈련 비용을 피하면서, 추론시에만 언어 모델 "
        "오케스트레이션을 바꿔 성능·해석가능성을 분리 평가할 수 있다.")
li(tf3, "동시에 CoT를 통하여 Long-Horizon 태스크 성능 향상을 기대해 볼 수 있다.")


# ══════════════════════════════════════════════════════════
# Slide 4 ── 2. 문제 정의 및 연구 목표
# ══════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_hdr(s4, "2. 문제 정의 및 연구 목표", 4); add_ftr(s4)
tf4 = new_tf(s4)

h1(tf4, "해결하고자 하는 문제", sp=2)
li(tf4, "고정된 UVA 정책·순동역학을 유지한 채, 외부 LLM의 구조화된 CoT가 행동 선택·롤아웃 "
        "평가를 개선할 수 있는지, 그리고 순동역 예측과의 일관성까지 검증할 수 있는가.")

h1(tf4, "연구 목표")
li(tf4, "추론 파이프라인(지시문 → CoT 단계 → UVA 액션/순동역 예측)을 정의하고, "
        "벤치마크에서 성공률·장기 완료·예측 일관성 지표로 UVA 단독 대비 이득을 정량화한다.")

h1(tf4, "연구 질문 (Research Questions)")
li(tf4, "RQ1: 동일 UVA 정책에서 CoT 계획 유무가 성공률·장기 과제 완료에 성능향상을 미치는가?")
li(tf4, "RQ2: CoT가 순동역학으로 생성한 미래 관측과 더 잘 맞는(일관된) 궤적·행동을 고르는가?")

h1(tf4, "연구의 범위(Scope)")
li(tf4, "시뮬레이션 벤치마크(예: LIBERO 소수 태스크 또는 PushT 계열), 이미 학습이 완료된 "
        "UVA만을 사용, 추론 시 CoT 오케스트레이션 및 프롬프트 변형.")
li(tf4, "실제 로봇 배포는 하지 않고 시뮬레이션만 사용.")


# ══════════════════════════════════════════════════════════
# Slide 5 ── 3. 관련 연구
# ══════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_hdr(s5, "3. 관련 연구 (Related Work)", 5); add_ftr(s5)
tf5 = new_tf(s5)

h1(tf5, "UVA 핵심 요약", sp=2)
li(tf5, "joint video-action latent + decoupled diffusion head를 통해 성능과 속도 균형 달성.")
li(tf5, "masking 기반 학습으로 정책/비디오/동역학 모델링을 통합.")
li(tf5, "PushT/Libero 등에서 strong multi-task 성능 보고.")

h1(tf5, "CoT 논문 핵심 요약")
li(tf5, "CoT는 중간 reasoning step을 통해 복잡한 문제를 단계적으로 해결.")
li(tf5, "최근 embodied robotics에서도 고수준 계획과 저수준 제어 연결을 위한 "
        "추론 구조의 중요성이 증가.")
li(tf5, "핵심은 \"중간 상태/목표를 명시적으로 모델링\"하는 것.")

h1(tf5, "기존 연구의 한계")
li(tf5, "UVA: 강한 표현력에도 불구하고 명시적 추론 체인이 약해 장기 계획에서 한계 가능.")
li(tf5, "일반 CoT 연구: 텍스트/추론 벤치 중심이 많아, 연속 제어 행동 생성으로의 직접 연결이 쉽지 않음.")
li(tf5, "따라서 \"비디오-액션 생성 모델 내부로 CoT를 실용적으로 주입하는 방식\"이 "
        "아직 충분히 정립되지 않음.")

h1(tf5, "차별점")
li(tf5, "본 연구는 UVA를 직접 확장해 실행 가능한 CoT 통합 정책을 제시.")
li(tf5, "단순 성능 비교를 넘어 \"중간 추론 품질(reasoning accuracy)\"까지 정량화.")
li(tf5, "단일 태스크 성능 보존 + 장기 복합 태스크 개선의 균형을 목표로 함.")


# ══════════════════════════════════════════════════════════
# Slide 6 ── 4. 제안 방법 (전체 아키텍처)
# ══════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_hdr(s6, "4. 제안 방법 (UVA + CoT 통합)", 6); add_ftr(s6)

tf6 = new_tf(s6, h=Inches(0.65))
h1(tf6, "전체 아키텍처(Overview)", sp=2)

ph_box(s6, Inches(0.4), Inches(1.55), W - Inches(0.8), Inches(3.0),
       "[ Agent Architecture: Input-Reasoning-Engine 다이어그램 이미지 삽입 ]")

tf6b = new_tf(s6, top=Inches(4.65), h=Inches(2.3))
h1(tf6b, "흐름:", sp=2)
li(tf6b, "1. LLM이 CoT로 \"지금 단계의 부분목표\"만 출력 (예: \"서랍 손잡이에 그리퍼 정렬\").")
li(tf6b, "2. 그 문자열을 language_goal로 바꿔 UVA가 짧은 구간 액션 청크를 예측·실행.")
li(tf6b, "3. 몇 스텝 후 관측이 바뀌면 다시 LLM이 다음 부분목표를 갱신.")

p_note = tf6b.add_paragraph()
p_note.space_before = Pt(4)
r_n = p_note.add_run()
r_n.text = "CLIP을 통해 CoT의 각 부분 목표를 벡터화시켜서 1차원으로 펼친 비디오 토큰 앞에 concat함."
r_n.font.size = Pt(12); r_n.font.name = FONT_KO; r_n.font.color.rgb = C_BODY


# ══════════════════════════════════════════════════════════
# Slide 7 ── 4. 제안 방법 (상세)
# ══════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_hdr(s7, "4. 제안 방법 (UVA + CoT 통합)", 7); add_ftr(s7)
tf7 = new_tf(s7)

h1(tf7, "핵심 아이디어", sp=2)
li(tf7, "기존에 Long-Horizon 태스크의 instruction 텍스트를 CLIP에 넣어 벡터로 바꾼 후 이미지 "
        "토큰에 접속(concat)하는 알고리즘이 있다. 이는 언뜻 이미지만 봐서는 태스크의 목적을 "
        "알 수 없던 Long-Horizon 태스크를 텍스트 instruction을 CLIP을 통해 따로 처리하여 로봇이 "
        "태스크 목표를 쫓을 수 있었다. 이 코드를 재활용하여 생성된 CoT의 각 단계를 역시 CLIP에 "
        "넣고 부분태스크의 경계에 해당하는 이미지에 접속(concat)하는 방식으로 CoT를 구현하는 "
        "것이 저의 아이디어입니다.")
li(tf7, "'현재 관측에서 무엇이 중요한가' → '다음에 만족해야 할 조건' → "
        "'후보 행동/롤아웃 비교' → UVA로 실행·예측.")

h1(tf7, "입력 및 출력 정의")
li(tf7, "Input: 이미지, 언어 instruction, 짧은 history(선택).")
li(tf7, "Output: CoT 중간 문자열(선택적으로 로그), 최종 액션 또는 "
        "UVA 순동역에 의한 미래 프레임/미래 액션")

h1(tf7, "Loss 및 학습 전략")
li(tf7, "본 연구 단계에서는 UVA 추가 학습 없음. 손실 최적화는 UVA 체크포인트에 적용하지 않는다. "
        "선택 시 LLM만 소량 in-context 예시 또는 경량 파인튜닝으로 프롬프트 적응 범위를 명시한다.")


# ══════════════════════════════════════════════════════════
# Slide 8 ── 5. 실험 설계
# ══════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_hdr(s8, "5. 실험 설계 (Experiment Design)", 8); add_ftr(s8)
tf8 = new_tf(s8)

h1(tf8, "데이터셋", sp=2)
li(tf8, "LIBERO(Long-Horizon Benchmark) 또는 PushT 계열 시뮬레이션")
li(tf8, "선정 이유: Long-Horizon 태스크·성공 여부가 명확하고, 시각 관측 기반 정책과 "
        "CoT 설명을 연결하기 용이함. 또한 Libero 10과 PushT계열은 실물 로봇 없이 "
        "시뮬레이션으로 돌릴 수 있는 데이터 셋임.")

h1(tf8, "Baseline")
li(tf8, "UVA(원본)")
li(tf8, "UVA + Zero-Shot 프롬프트")

h1(tf8, "평가 지표(Metrics)")
li(tf8, "Success rate(및 필요 시 completion rate).")
li(tf8, "선택 지표: CoT가 고른 액션이 UVA 예측과 얼마나 일관적인지 "
        "(예: 예측 프레임·상태와의 정합 스코어)")


# ══════════════════════════════════════════════════════════
# Slide 9 ── 6. 기대 효과 및 활용 방안
# ══════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
add_hdr(s9, "6. 기대 효과 및 활용 방안", 9); add_ftr(s9)
tf9 = new_tf(s9)

h1(tf9, "학술적 기여", sp=2)
li(tf9, "UVA 모델 위에 언어 CoT 오케스트레이션만 바꿔도 장기 과제와 예측 일관성이 달라질 수 "
        "있음을 실증하고, 베이스라인 축을 명확히 제시한다.")

h1(tf9, "실용적 활용 방안", sp=8)
li(tf9, "가정·물류 등 지시 기반 조작에서 사람이 읽을 수 있는 계획 단계와 검증 가능한 "
        "동역학 예측을 함께 제공하는 의사결정 스택으로 확장 가능하다.")

h1(tf9, "후속 연구 방향", sp=8)
li(tf9, "실제 로봇·멀티모달 관측 확장, CoT 단계와 UVA 내부 표현의 정렬, "
        "제한적 LLM 적응과의 비교.")


# ══════════════════════════════════════════════════════════
OUT = r"c:\Users\werty\Documents\GitHub\unified_video_action\UVA_to_CoT_기획서.pptx"
prs.save(OUT)
print("Saved:", OUT)
